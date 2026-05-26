#!/usr/bin/env python
"""家长会PPT — 高级美化版"""
import io, os, secrets
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw, ImageFilter

# ═══════════════════════════════════ 常量 ═══════════════════════════════════
DEEP_BLUE = RGBColor(0x1B, 0x3A, 0x4B)
WARM_GOLD = RGBColor(0xC6, 0x9C, 0x6D)
CREAM_WHITE = RGBColor(0xFB, 0xF7, 0xF0)
CORAL_RED = RGBColor(0xE0, 0x7A, 0x5F)
LIGHT_GRAY = RGBColor(0xD5, 0xD5, 0xD5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x88, 0x88, 0x88)
PALE_BLUE = RGBColor(0xE8, 0xEE, 0xF2)

TITLE_FONT = "思源宋体 Heavy"
BODY_FONT = "思源黑体 Regular"
DATA_FONT = "DIN Alternate Bold"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(1.07)   # ~8%
MARGIN_R = Inches(1.07)
MARGIN_TOP = Inches(0.9)  # ~12%

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ═══════════════════════════════════ 纹理生成 ═══════════════════════════════════
def make_texture_tile():
    """生成120x120的网格+纸纹拼接图，RGBA PNG"""
    size = 120
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # 40px间距网格线，深青蓝 5%透明度
    for x in range(0, size, 40):
        draw.line([(x, 0), (x, size)], fill=(0x1B, 0x3A, 0x4B, 13))
    for y in range(0, size, 40):
        draw.line([(0, y), (size, y)], fill=(0x1B, 0x3A, 0x4B, 13))
    # 纤维噪点 3%透明度
    for _ in range(300):
        x = secrets.randbelow(size)
        y = secrets.randbelow(size)
        a = secrets.randbelow(8)
        if a > 0:
            draw.point((x, y), fill=(0x1B, 0x3A, 0x4B, a))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

texture_png = make_texture_tile()
texture_dir = "E:/git/test_app/_texture_tile.png"
with open(texture_dir, "wb") as f:
    f.write(texture_png)

# ═══════════════════════════════════ 工具函数 ═══════════════════════════════════
def set_bg(slide, color=CREAM_WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_texture(slide):
    """叠加纹理层"""
    pic = slide.shapes.add_picture(texture_dir, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # 设为图片平铺
    blipFill = pic._element.findall(".//" + qn("a:blipFill"))[0] if pic._element.findall(".//" + qn("a:blipFill")) else None
    if blipFill is None:
        blipFill = pic._element.find(qn("p:blipFill"))
    if blipFill is not None:
        stretch = blipFill.find(qn("a:stretch"))
        if stretch is not None:
            blipFill.remove(stretch)
        tile = etree.SubElement(blipFill, qn("a:tile"))
        tile.set("tx", "0")
        tile.set("ty", "0")
        tile.set("sx", "100%")
        tile.set("sy", "100%")
        tile.set("flip", "none")
        tile.set("algn", "tl")

def add_shadow(shape, blur_pt=8, dist_pt=2, angle=135, opacity_pct=15):
    """给shape加柔投影"""
    spPr = get_spPr(shape)
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", str(Emu(Pt(blur_pt))))
    outerShdw.set("dist", str(Emu(Pt(dist_pt))))
    outerShdw.set("dir", str(angle * 60000))
    outerShdw.set("algn", "tl")
    srgb = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgb.set("val", "1B3A4B")
    alpha_elem = etree.SubElement(srgb, qn("a:alpha"))
    alpha_elem.set("val", str(opacity_pct * 1000))

def get_spPr(shape):
    """获取shape的spPr元素"""
    for ns in ("a:spPr", "p:spPr"):
        sp = shape._element.find(qn(ns))
        if sp is not None:
            return sp
    return shape._element

def add_gradient_fill(shape, color1, color2, angle_deg=0, path="linear"):
    """渐变填充"""
    spPr = get_spPr(shape)
    # Remove existing solid fill
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    for old in spPr.findall(qn("a:noFill")):
        spPr.remove(old)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    if path == "linear":
        lin = etree.SubElement(gradFill, qn("a:lin"))
        lin.set("ang", str(angle_deg * 60000))
        lin.set("scaled", "0")

    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, color in [(0, color1), (100000, color2)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")

def add_radial_gradient(shape, center_color, edge_color, center_alpha=12, edge_alpha=0):
    """径向渐变"""
    spPr = get_spPr(shape)
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    for old in spPr.findall(qn("a:noFill")):
        spPr.remove(old)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    path_elem = etree.SubElement(gradFill, qn("a:path"))
    path_elem.set("path", "circle")
    fillRect = etree.SubElement(path_elem, qn("a:fillToRect"))
    fillRect.set("l", "50000")
    fillRect.set("t", "50000")
    fillRect.set("r", "50000")
    fillRect.set("b", "50000")

    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, color, alpha in [(0, center_color, center_alpha), (100000, edge_color, edge_alpha)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
        if alpha > 0:
            alpha_elem = etree.SubElement(srgb, qn("a:alpha"))
            alpha_elem.set("val", str(alpha * 1000))

def add_textbox(slide, left, top, width, height, text="", font_name=BODY_FONT,
                font_size=14, color=DEEP_BLUE, bold=False, alignment=PP_ALIGN.LEFT,
                line_spacing=1.5, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines):
    """lines = [(text, font_name, font_size, color, bold, alignment, ls), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, fn, fs, col, bld, aln, ls) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = fn
        p.font.size = Pt(fs)
        p.font.color.rgb = col
        p.font.bold = bld
        p.alignment = aln
        p.line_spacing = Pt(fs * ls)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius / max(width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def add_circle(slide, left, top, diameter, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=WARM_GOLD, width=Pt(1)):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    connector.line.fill.solid()
    return connector


# ═══════════════════════════════════ 第1页：封面 ═══════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, CREAM_WHITE)
add_texture(slide)

# 右下角径向渐变光晕
glow_w = Inches(5.5)
glow_h = Inches(5.5)
glow_shape = add_circle(slide, SLIDE_W - glow_w - Inches(1.5), SLIDE_H - glow_h - Inches(0.8),
                         glow_w, DEEP_BLUE)
add_radial_gradient(glow_shape, WARM_GOLD, WARM_GOLD, center_alpha=12, edge_alpha=0)

# 装饰曲线 — 使用自由曲线加渐变填充
# 简化：创建一条自由形曲线，从右下光晕蜿蜒向上
from pptx.oxml.ns import qn as qn_xml
curve_xml = """
<p:cxnSp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:p="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvCxnSpPr>
    <p:cNvPr id="0" name="Curve"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="0" y="0"/>
      <a:ext cx="1" cy="1"/>
    </a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:pathLst>
        <a:path w="1" h="1" fill="none">
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          <a:cubicBezTo>
            <a:pt x="r" y="0"/><a:pt x="r" y="1"/><a:pt x="r" y="1"/>
          </a:cubicBezTo>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:ln w="25400" cap="rnd">
      <a:solidFill><a:srgbClr val="C69C6D"/></a:solidFill>
    </a:ln>
  </p:spPr>
</p:cxnSp>"""

# 简化装饰曲线改用多个线段+圆点
curve_start_x = Inches(10.0)
curve_start_y = Inches(5.5)
pts = [
    (curve_start_x, Inches(5.5)),
    (Inches(10.5), Inches(4.8)),
    (Inches(10.2), Inches(4.0)),
    (Inches(9.8), Inches(3.5)),
    (Inches(9.2), Inches(3.2)),
]
# 逐段画曲线，宽度递减
widths = [Pt(6), Pt(4.5), Pt(3), Pt(2)]
for i in range(len(pts) - 1):
    add_line(slide, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], DEEP_BLUE, widths[i])

# 末端3个暖金圆点
dot_x, dot_y = Inches(9.0), Inches(3.0)
dot_sizes = [Inches(0.12), Inches(0.08), Inches(0.05)]
dot_positions = [(0, 0), (Inches(0.3), Inches(-0.15)), (Inches(0.15), Inches(-0.35))]
for ds, (dx, dy) in zip(dot_sizes, dot_positions):
    add_circle(slide, dot_x + dx, dot_y + dy, ds, WARM_GOLD)

# 主标题（加文字投影效果）
title_box = add_textbox(slide, Inches(1.2), Inches(1.3), Inches(10.5), Inches(2.0),
                        "2025-2026学年第二学期\n期中家长会",
                        font_name=TITLE_FONT, font_size=52, color=DEEP_BLUE, bold=True,
                        alignment=PP_ALIGN.LEFT, line_spacing=1.25)
# 标题阴影（通过XML添加）
spPr = title_box._element.find(qn("p:txBody"))
# 给文字加投影 — 在段落属性层面
for p in title_box.text_frame.paragraphs:
    pPr = p._pPr if p._pPr is not None else etree.SubElement(title_box.text_frame._txBody, qn("a:pPr"))
    # 文字阴影通过 defRPr 设置
    for run in p.runs:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(run._r, qn("a:rPr"))
        effectLst_r = rPr.find(qn("a:effectLst"))
        if effectLst_r is None:
            effectLst_r = etree.SubElement(rPr, qn("a:effectLst"))
        outerShdw = etree.SubElement(effectLst_r, qn("a:outerShdw"))
        outerShdw.set("blurRad", "25400")
        outerShdw.set("dist", "25400")
        outerShdw.set("dir", "2700000")
        outerShdw.set("algn", "tl")
        srgb = etree.SubElement(outerShdw, qn("a:srgbClr"))
        srgb.set("val", "C69C6D")
        alpha_elem = etree.SubElement(srgb, qn("a:alpha"))
        alpha_elem.set("val", "20000")

# 副标题上方40px细线
add_line(slide, Inches(1.2), Inches(3.3), Inches(2.2), Inches(3.3), WARM_GOLD, Pt(1.5))
# 副标题
add_textbox(slide, Inches(1.2), Inches(3.4), Inches(10), Inches(0.7),
            "同 心 同 行 ， 看 见 成 长",
            font_name=TITLE_FONT, font_size=24, color=WARM_GOLD, bold=False,
            alignment=PP_ALIGN.LEFT, line_spacing=1.5)

# 暖金渐变分隔线
sep_line = add_line(slide, Inches(1.2), Inches(4.3), Inches(5.5), Inches(4.3), WARM_GOLD, Pt(0.8))

# 学校/班级信息
add_multiline_textbox(slide, Inches(1.2), Inches(4.5), Inches(5), Inches(2.0), [
    ("某某中学", BODY_FONT, 14, GRAY_TEXT, False, PP_ALIGN.LEFT, 1.8),
    ("七年级（3）班", BODY_FONT, 14, GRAY_TEXT, False, PP_ALIGN.LEFT, 1.8),
    ("班主任：张老师  ·  2026年5月", BODY_FONT, 14, GRAY_TEXT, False, PP_ALIGN.LEFT, 1.8),
])


# ═══════════════════════════════════ 第2页：会议议程 ═══════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, CREAM_WHITE)
add_texture(slide)

# 左侧35%渐变底色
left_bg = add_rect(slide, Inches(0), Inches(0), Inches(4.7), SLIDE_H, fill_color=CREAM_WHITE)
add_gradient_fill(left_bg, DEEP_BLUE, CREAM_WHITE, angle_deg=0, path="linear")
# Override alpha in gradient stops for transparency
spPr = get_spPr(left_bg)
gradFill = spPr.find(qn("a:gradFill"))
gsLst = gradFill.find(qn("a:gsLst"))
for gs in gsLst.findall(qn("a:gs")):
    srgb = gs.find(qn("a:srgbClr"))
    if srgb is not None:
        existing_alpha = srgb.find(qn("a:alpha"))
        if existing_alpha is None:
            alpha_elem = etree.SubElement(srgb, qn("a:alpha"))
            pos = int(gs.get("pos"))
            if pos == 0:
                alpha_elem.set("val", str(8 * 1000))
            else:
                alpha_elem.set("val", "0")

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
            "会议议程", font_name=TITLE_FONT, font_size=28, color=DEEP_BLUE, bold=True)

agenda = [
    ("壹", "学期回顾与数据分析", "期中考试成绩总览、班级亮点与不足", False),
    ("贰", "学习表现深度解读", "各科进步情况、课堂表现评估", True),
    ("叁", "共性问题与配合建议", "手机管理、作业习惯、心理状态", False),
    ("肆", "家校沟通与答疑", "自由交流、个别咨询安排", False),
]

tl_left = Inches(2.2)
tl_top = Inches(1.6)
node_d = Inches(0.16)
spacing = Inches(1.3)
card_left = Inches(4.2)
card_w = Inches(7.8)
card_h = Inches(1.05)

# 时间轴线
add_line(slide, tl_left + node_d // 2, tl_top, tl_left + node_d // 2,
         tl_top + spacing * (len(agenda) - 1), DEEP_BLUE, Pt(1.2))

for i, (num, title, desc, is_current) in enumerate(agenda):
    y = tl_top + spacing * i

    # 光晕圆环
    ring = add_circle(slide, tl_left - Inches(0.04), y - Inches(0.04),
                      node_d + Inches(0.08), WARM_GOLD)
    ring.fill.background()
    ring.line.color.rgb = WARM_GOLD
    ring.line.width = Pt(1.5)
    ring.line.fill.solid()
    # 给圆环加透明度
    spPr_r = ring._element.find(qn("p:spPr"))
    ln = spPr_r.find(qn("a:ln"))
    if ln is not None:
        solidFill = ln.find(qn("a:solidFill"))
        if solidFill is not None:
            srgb = solidFill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn("a:alpha"))
                alpha_elem.set("val", "30000")

    # 实心节点
    dot_color = CORAL_RED if is_current else WARM_GOLD
    add_circle(slide, tl_left, y, node_d, dot_color)

    # 编号
    add_textbox(slide, tl_left + Inches(0.35), y - Inches(0.02), Inches(0.5), Inches(0.25),
                num, font_name=TITLE_FONT, font_size=14, color=DEEP_BLUE, bold=True)

    # 卡片
    border_c = CORAL_RED if is_current else DEEP_BLUE
    border_w = Pt(4) if is_current else Pt(1.5)
    card = add_rect(slide, card_left, y - Inches(0.1), card_w, card_h,
                    fill_color=WHITE, line_color=border_c, line_width=border_w)
    add_shadow(card, blur_pt=6, dist_pt=2, angle=135, opacity_pct=12)

    # 当前项右上角红点指示灯
    if is_current:
        dot = add_circle(slide, card_left + card_w - Inches(0.35), y - Inches(0.03),
                         Inches(0.1), CORAL_RED)
        # 加呼吸光晕
        spPr_dot = dot._element.find(qn("p:spPr"))
        effectLst_d = spPr_dot.find(qn("a:effectLst"))
        if effectLst_d is None:
            effectLst_d = etree.SubElement(spPr_dot, qn("a:effectLst"))
        glow = etree.SubElement(effectLst_d, qn("a:glow"))
        glow.set("rad", str(Emu(Pt(4))))
        srgb_g = etree.SubElement(glow, qn("a:srgbClr"))
        srgb_g.set("val", "E07A5F")
        alpha_g = etree.SubElement(srgb_g, qn("a:alpha"))
        alpha_g.set("val", "50000")

    # 标题
    add_textbox(slide, card_left + Inches(0.25), y + Inches(0.05),
                card_w - Inches(0.5), Inches(0.3),
                title, font_name=BODY_FONT, font_size=18, color=DEEP_BLUE, bold=True)
    # 暖金分隔线
    add_line(slide, card_left + Inches(0.25), y + Inches(0.38),
             card_left + Inches(1.5), y + Inches(0.38), WARM_GOLD, Pt(0.5))
    # 描述
    add_textbox(slide, card_left + Inches(0.25), y + Inches(0.46),
                card_w - Inches(0.5), Inches(0.4),
                desc, font_name=BODY_FONT, font_size=14, color=GRAY_TEXT, bold=False)


# ═══════════════════════════════════ 第3页：学习表现 ═══════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, CREAM_WHITE)
add_texture(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
            "学习表现数据：期中成绩对比", font_name=TITLE_FONT, font_size=28,
            color=DEEP_BLUE, bold=True)

chart_data = CategoryChartData()
chart_data.categories = ['语文', '数学', '英语', '物理', '历史']
chart_data.add_series('班级平均', (82, 78, 71, 85, 88))
chart_data.add_series('年级平均', (80, 83, 76, 80, 86))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.2), Inches(1.3), Inches(9.0), Inches(4.0),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
chart.legend.font.name = BODY_FONT

plot = chart.plots[0]
plot.gap_width = 80

s1 = plot.series[0]
s1.format.fill.solid()
s1.format.fill.fore_color.rgb = DEEP_BLUE

s2 = plot.series[1]
s2.format.fill.solid()
s2.format.fill.fore_color.rgb = LIGHT_GRAY

value_axis = chart.value_axis
value_axis.minimum_scale = 50
value_axis.maximum_scale = 100
value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
value_axis.format.line.fill.background()

cat_axis = chart.category_axis
cat_axis.format.line.fill.background()
cat_axis.tick_labels.font.name = BODY_FONT
cat_axis.tick_labels.font.size = Pt(12)

# 图表区域加浅投影
add_shadow(chart_frame, blur_pt=10, dist_pt=3, angle=135, opacity_pct=10)

# 结论区
add_textbox(slide, Inches(1.2), Inches(5.5), Inches(10.5), Inches(0.4),
            "▇ 数学（-5分）年级差距最大       ▇ 物理（+5分）班级优势明显       ▇ 英语（-5分）需重点关注",
            font_name=BODY_FONT, font_size=15, color=DEEP_BLUE, bold=True)

add_textbox(slide, Inches(1.2), Inches(5.95), Inches(10.5), Inches(0.4),
            "数学和英语是唯二低于或等于年级平均的学科，建议优先强化这两个科目的基础训练与错题整理。",
            font_name=BODY_FONT, font_size=13, color=GRAY_TEXT, bold=False)


# ═══════════════════════════════════ 第4页：共性问题 ═══════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, CREAM_WHITE)
add_texture(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
            "共性问题与配合建议", font_name=TITLE_FONT, font_size=28, color=DEEP_BLUE, bold=True)

problems = [
    ("1", "手机使用时间过长", "→ 执行手机契约\n→ 检查过程痕迹"),
    ("2", "作业完成质量不高", "→ 检查过程痕迹\n→ 建立错题本制度"),
    ("3", "课堂专注度下滑", "→ 每日3分钟复盘\n→ 保证充足睡眠"),
]

for i, (num, title, actions) in enumerate(problems):
    y = Inches(1.3) + Inches(1.9) * i

    # 左侧问题卡片（渐变背景）
    card_left = Inches(0.8)
    card_w = Inches(4.5)
    card_h = Inches(1.6)
    card = add_rect(slide, card_left, y, card_w, card_h, fill_color=PALE_BLUE,
                    line_color=DEEP_BLUE, line_width=Pt(0.5))
    add_shadow(card, blur_pt=5, dist_pt=1.5, angle=135, opacity_pct=10)

    # 珊瑚红数字 + 圆形底色
    num_circle = add_circle(slide, Inches(1.1), y + Inches(0.15), Inches(0.55), CORAL_RED)
    # 数字下方透明度10%圆
    num_circle2 = add_circle(slide, Inches(1.05), y + Inches(0.1), Inches(0.65), CORAL_RED)
    spPr_nc = num_circle2._element.find(qn("p:spPr"))
    solid_nc = spPr_nc.find(qn("a:solidFill"))
    if solid_nc is not None:
        srgb_nc = solid_nc.find(qn("a:srgbClr"))
        if srgb_nc is not None:
            alpha_nc = etree.SubElement(srgb_nc, qn("a:alpha"))
            alpha_nc.set("val", "10000")

    add_textbox(slide, Inches(1.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                num, font_name=DATA_FONT, font_size=30, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 问题标题
    add_textbox(slide, Inches(2.0), y + Inches(0.2), Inches(3.0), Inches(0.35),
                title, font_name=BODY_FONT, font_size=17, color=DEEP_BLUE, bold=True)

    # 右侧行动建议
    sugg_left = Inches(5.8)
    sugg_w = Inches(6.5)

    # 图标圆形容器
    icon_circle = add_circle(slide, sugg_left, y + Inches(0.15), Inches(0.45), WARM_GOLD)
    icon_circle.fill.background()
    icon_circle.line.color.rgb = WARM_GOLD
    icon_circle.line.width = Pt(1.5)
    icon_circle.line.fill.solid()
    spPr_ic = icon_circle._element.find(qn("p:spPr"))
    ln_ic = spPr_ic.find(qn("a:ln"))
    if ln_ic is not None:
        sf_ic = ln_ic.find(qn("a:solidFill"))
        if sf_ic is not None:
            srgb_ic = sf_ic.find(qn("a:srgbClr"))
            if srgb_ic is not None:
                alpha_ic = etree.SubElement(srgb_ic, qn("a:alpha"))
                alpha_ic.set("val", "40000")

    add_textbox(slide, sugg_left + Inches(0.12), y + Inches(0.2), Inches(0.3), Inches(0.3),
                "→", font_name=BODY_FONT, font_size=14, color=DEEP_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, sugg_left + Inches(0.6), y + Inches(0.15), Inches(5.5), Inches(1.2),
                actions, font_name=BODY_FONT, font_size=15, color=DEEP_BLUE, bold=True,
                line_spacing=1.8)

    # 浅灰虚线分隔
    if i < 2:
        sep_y = y + Inches(1.75)
        sep_line = add_line(slide, Inches(0.8), sep_y, Inches(12.5), sep_y,
                            RGBColor(0xCC, 0xC8, 0xC0), Pt(0.5))


# ═══════════════════════════════════ 第5页：家校沟通 ═══════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, CREAM_WHITE)

# 中心射线光晕
glow_center = add_circle(slide, Inches(3.0), Inches(-1.5), Inches(7.0), WARM_GOLD)
add_radial_gradient(glow_center, WARM_GOLD, WARM_GOLD, center_alpha=12, edge_alpha=0)

# 纹理叠在最上层
add_texture(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
            "家校沟通", font_name=TITLE_FONT, font_size=28, color=DEEP_BLUE, bold=True)

# 照片占位区 — 圆角+投影
photo = add_rect(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(3.8),
                 fill_color=RGBColor(0xE8, 0xE3, 0xD8),
                 line_color=WHITE, line_width=Pt(3), corner_radius=Inches(0.22))
add_shadow(photo, blur_pt=12, dist_pt=4, angle=135, opacity_pct=18)

# 相机图标占位
add_textbox(slide, Inches(4.5), Inches(2.7), Inches(5.0), Inches(0.8),
            "📷  插入班级合照", font_name=BODY_FONT, font_size=20,
            color=GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# 联系方式
add_textbox(slide, Inches(0.8), Inches(5.3), Inches(4.5), Inches(0.35),
            "联系方式", font_name=BODY_FONT, font_size=15, color=WARM_GOLD, bold=True)

add_multiline_textbox(slide, Inches(0.8), Inches(5.75), Inches(5.0), Inches(1.2), [
    ("班主任电话：138-xxxx-xxxx", BODY_FONT, 13, GRAY_TEXT, False, PP_ALIGN.LEFT, 2.2),
    ("办公时间：周一至周五 8:00-17:00", BODY_FONT, 13, GRAY_TEXT, False, PP_ALIGN.LEFT, 2.2),
    ("班级通知群：微信群 / 钉钉群", BODY_FONT, 13, GRAY_TEXT, False, PP_ALIGN.LEFT, 2.2),
])

# 装饰引号（暖金，透明度10%）
add_textbox(slide, Inches(6.8), Inches(5.1), Inches(2.0), Inches(1.3),
            "❝", font_name=TITLE_FONT, font_size=80, color=WARM_GOLD, bold=False,
            alignment=PP_ALIGN.LEFT, line_spacing=1.0)
# 给引号加透明度
# Note: text transparency in PPTX is complex, using the color itself as approximation

# 结束语 — 字号第二大
add_textbox(slide, Inches(7.0), Inches(5.3), Inches(5.5), Inches(1.2),
            "没有完美的孩子，\n也没有完美的家长，\n只有不断靠近真实的理解。",
            font_name=TITLE_FONT, font_size=22, color=DEEP_BLUE, bold=False,
            alignment=PP_ALIGN.CENTER, line_spacing=2.0)

# 底部细线
add_line(slide, Inches(0.8), Inches(7.0), Inches(12.5), Inches(7.0), WARM_GOLD, Pt(0.5))


# ═══════════════════════════════════ 保存 ═══════════════════════════════════
output_path = "E:/git/test_app/家长会PPT.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
