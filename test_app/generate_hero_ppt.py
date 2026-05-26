#!/usr/bin/env python
"""英雄主题家长会PPT — 抗战纪念版"""
import io, os, secrets, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw, ImageFilter, ImageFont

# ═══════════════════ 配色 ═══════════════════
HERO_RED = RGBColor(0xC4, 0x1E, 0x3A)
STAR_GOLD = RGBColor(0xD4, 0xA8, 0x43)
SMOKE_GRAY = RGBColor(0xF2, 0xEC, 0xE4)
CAST_IRON = RGBColor(0x2C, 0x2C, 0x2C)
VICTORY_BLUE = RGBColor(0x2B, 0x5B, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x88, 0x88, 0x88)

TITLE_FONT = "演示悠然小楷"
BODY_FONT = "阿里巴巴普惠体 Medium"
DATA_FONT = "DIN Condensed Bold"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M_LR = Inches(1.33)  # ≥10%
M_TB = Inches(0.75)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ═══════════════════ 纹理生成 ═══════════════════
def make_manuscript_texture():
    """义勇军进行曲手稿纹理 — 12%透明度"""
    w, h = 400, 300
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # 五线谱横线
    for i in range(25):
        y = 10 + i * 12
        alpha = 31  # 12% of 255
        draw.line([(0, y), (w, y)], fill=(0x2C, 0x2C, 0x2C, alpha))
    # 随机音符/记号
    for _ in range(60):
        x = secrets.randbelow(w)
        y = secrets.randbelow(h)
        r = secrets.randbelow(4) + 1
        draw.ellipse([(x, y), (x + r * 2, y + r)], fill=(0x2C, 0x2C, 0x2C, 25))
    # 稀疏竖线（小节线）
    for x in range(30, w, 60):
        draw.line([(x, 5), (x, h - 5)], fill=(0x2C, 0x2C, 0x2C, 15))
    return img

def make_rays_texture():
    """金色放射光芒 — 6%透明度，中心放射"""
    size = 500
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    cx, cy = size // 2, size // 2
    for angle in range(0, 360, 12):
        rad = math.radians(angle)
        ex = cx + int(size * 0.7 * math.cos(rad))
        ey = cy + int(size * 0.7 * math.sin(rad))
        draw.line([(cx, cy), (ex, ey)], fill=(0xD4, 0xA8, 0x43, 15), width=1)
        # 较短的辅射线
        ex2 = cx + int(size * 0.35 * math.cos(rad + 0.05))
        ey2 = cy + int(size * 0.35 * math.sin(rad + 0.05))
        draw.line([(cx + 5, cy + 5), (ex2, ey2)], fill=(0xD4, 0xA8, 0x43, 10), width=1)
    return img

def save_tile(img, path):
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    with open(path, "wb") as f:
        f.write(buf.getvalue())

manuscript_img = make_manuscript_texture()
rays_img = make_rays_texture()
manu_path = "E:/git/test_app/_manu_tile.png"
rays_path = "E:/git/test_app/_rays_tile.png"
save_tile(manuscript_img, manu_path)
save_tile(rays_img, rays_path)


# ═══════════════════ 工具函数 ═══════════════════
def get_spPr(shape):
    for ns in ("a:spPr", "p:spPr"):
        sp = shape._element.find(qn(ns))
        if sp is not None:
            return sp
    return shape._element

def set_bg(slide, color=SMOKE_GRAY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_texture_tile(slide, path, alpha=None):
    """全幅平铺纹理"""
    pic = slide.shapes.add_picture(path, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    blipFill = pic._element.findall(".//" + qn("a:blipFill"))[0] if pic._element.findall(".//" + qn("a:blipFill")) else pic._element.find(qn("p:blipFill"))
    if blipFill is not None:
        stretch = blipFill.find(qn("a:stretch"))
        if stretch is not None:
            blipFill.remove(stretch)
        tile = etree.SubElement(blipFill, qn("a:tile"))
        tile.set("tx", "0"); tile.set("ty", "0")
        tile.set("sx", "100%"); tile.set("sy", "100%")
        tile.set("flip", "none"); tile.set("algn", "tl")
    return pic

def add_shadow(shape, blur_pt=6, dist_pt=2, angle=135, opacity_pct=15):
    spPr = get_spPr(shape)
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", str(Emu(Pt(blur_pt))))
    outerShdw.set("dist", str(Emu(Pt(dist_pt))))
    outerShdw.set("dir", str(angle * 60000))
    outerShdw.set("algn", "tl")
    srgb = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgb.set("val", "2C2C2C")
    alpha_e = etree.SubElement(srgb, qn("a:alpha"))
    alpha_e.set("val", str(opacity_pct * 1000))

def add_radial_glow(shape, color, center_alpha=50, edge_alpha=0):
    spPr = get_spPr(shape)
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    for old in spPr.findall(qn("a:noFill")):
        spPr.remove(old)
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    path_e = etree.SubElement(gradFill, qn("a:path"))
    path_e.set("path", "circle")
    fillRect = etree.SubElement(path_e, qn("a:fillToRect"))
    for a in ("l", "t", "r", "b"):
        fillRect.set(a, "50000")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, ca, ea in [(0, center_alpha, edge_alpha), (100000, edge_alpha, 0)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
        alpha_e = etree.SubElement(srgb, qn("a:alpha"))
        alpha_e.set("val", str(ca * 1000 if pos == 0 else ea * 1000))

def add_linear_gradient(shape, c1, c2, angle=0):
    spPr = get_spPr(shape)
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    for old in spPr.findall(qn("a:noFill")):
        spPr.remove(old)
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(angle * 60000)); lin.set("scaled", "0")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, c in [(0, c1), (100000, c2)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{c[0]:02X}{c[1]:02X}{c[2]:02X}")

def add_textbox(slide, left, top, width, height, text="", font_name=BODY_FONT,
                font_size=14, color=CAST_IRON, bold=False, alignment=PP_ALIGN.LEFT,
                line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text; p.font.name = font_name
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline(slide, left, top, width, height, lines):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, (txt, fn, fs, c, b, al, ls) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.name = fn; p.font.size = Pt(fs)
        p.font.color.rgb = c; p.font.bold = b; p.alignment = al
        p.line_spacing = Pt(fs * ls)
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.fill.solid()
        if lw: s.line.width = lw
    else:
        s.line.fill.background()
    return s

def add_line(slide, x1, y1, x2, y2, color, width=Pt(1)):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = width; c.line.fill.solid()
    return c

def add_star(slide, left, top, size, fill=STAR_GOLD):
    s = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, left, top, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    # 微调星形外观
    try: s.adjustments[0] = 0.4
    except: pass
    return s

def add_circle(slide, left, top, d, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.fill.solid()
    if lw: s.line.width = lw
    else: s.line.fill.background() if not line else None
    return s


# ═══════════════════ 生成白鸽剪影 PNG ═══════════════════
def make_dove_silhouette(size=80):
    """简易白鸽剪影"""
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # 简化鸽子轮廓
    body = [(size * 0.3, size * 0.55), (size * 0.15, size * 0.45), (size * 0.05, size * 0.35),
            (size * 0.1, size * 0.3), (size * 0.25, size * 0.25), (size * 0.4, size * 0.15),
            (size * 0.6, size * 0.05), (size * 0.8, size * 0.1), (size * 0.9, size * 0.2),
            (size * 0.85, size * 0.3), (size * 0.7, size * 0.35), (size * 0.6, size * 0.5),
            (size * 0.5, size * 0.55), (size * 0.35, size * 0.6)]
    draw.polygon(body, fill=(0xD4, 0xA8, 0x43, 180))
    # 翅膀
    wing = [(size * 0.45, size * 0.12), (size * 0.55, size * 0.02), (size * 0.7, size * 0.0),
            (size * 0.85, size * 0.08), (size * 0.75, size * 0.2), (size * 0.6, size * 0.18)]
    draw.polygon(wing, fill=(0xD4, 0xA8, 0x43, 160))
    buf = io.BytesIO(); img.save(buf, format="PNG"); return buf.getvalue()

dove_png = make_dove_silhouette()
dove_path = "E:/git/test_app/_dove.png"
with open(dove_path, "wb") as f:
    f.write(dove_png)


# ═══════════════════ 第1页：封面 ═══════════════════
slide = prs.slides.add_slide(blank)
set_bg(slide, SMOKE_GRAY)
add_texture_tile(slide, manu_path)
add_texture_tile(slide, rays_path)

# 太行山剪影 — 多边形
mountain_pts = [
    (Inches(0), SLIDE_H), (Inches(0), Inches(5.2)), (Inches(0.8), Inches(4.5)),
    (Inches(1.5), Inches(4.8)), (Inches(2.3), Inches(4.0)), (Inches(3.0), Inches(4.5)),
    (Inches(3.8), Inches(3.6)), (Inches(4.6), Inches(4.2)), (Inches(5.5), Inches(3.5)),
    (Inches(6.3), Inches(3.9)), (Inches(7.0), Inches(3.3)), (Inches(7.8), Inches(3.7)),
    (Inches(8.5), Inches(3.1)), (Inches(9.2), Inches(3.6)), (Inches(9.8), Inches(3.2)),
    (Inches(10.5), Inches(3.8)), (Inches(11.2), Inches(3.4)), (Inches(11.8), Inches(3.8)),
    (Inches(12.5), Inches(3.5)), (Inches(13.0), Inches(4.0)), (Inches(SLIDE_W), Inches(3.8)),
    (Inches(SLIDE_W), SLIDE_H),
]
# 使用 freeform
builder = slide.shapes.build_freeform(Emu(mountain_pts[0][0]), Emu(mountain_pts[0][1]))
for pt in mountain_pts[1:]:
    builder.add_line_segments([(Emu(pt[0]), Emu(pt[1]))], close=False)
mountain = builder.convert_to_shape()
mountain.fill.solid()
mountain.fill.fore_color.rgb = HERO_RED
mountain.line.fill.background()
# 渐变透明度
spPr_m = get_spPr(mountain)
solid_m = spPr_m.find(qn("a:solidFill"))
if solid_m is not None:
    srgb_m = solid_m.find(qn("a:srgbClr"))
    if srgb_m is not None:
        alpha_m = etree.SubElement(srgb_m, qn("a:alpha"))
        alpha_m.set("val", "40000")  # 60% opacity for mountain base

# 纪念碑轮廓 — 居中偏右
monu_x = Inches(8.5); monu_top = Inches(1.5)
monu_w = Inches(1.2); monu_h = Inches(3.0)
# 碑身
add_shape(slide, MSO_SHAPE.RECTANGLE, monu_x, monu_top + Inches(0.6), monu_w, monu_h - Inches(0.6),
          fill=STAR_GOLD)
add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, monu_x - Inches(0.2), monu_top,
          monu_w + Inches(0.4), Inches(0.7), fill=STAR_GOLD)
# 纪念碑光芒
monu_glow = add_circle(slide, monu_x - Inches(1.5), monu_top - Inches(1.0),
                       monu_w + Inches(3.5), fill=STAR_GOLD)
add_radial_glow(monu_glow, STAR_GOLD, center_alpha=35, edge_alpha=0)

# 左上角红星
star_size = Inches(1.0)
star = add_star(slide, Inches(0.6), Inches(0.4), star_size, STAR_GOLD)
add_shadow(star, blur_pt=8, dist_pt=3, angle=135, opacity_pct=25)

# 红绸飘带 — 使用波浪形矩形
ribbon = add_shape(slide, MSO_SHAPE.WAVE, Inches(1.5), Inches(0.9), Inches(10.5), Inches(0.5),
                   fill=HERO_RED)
add_linear_gradient(ribbon, HERO_RED, STAR_GOLD, angle=90)

# 主标题 — 竖排居左（模拟竖排：每个字一行）
title_chars = "以英雄之光照亮成长之路"
title_x = Inches(1.0); title_y = Inches(1.8)
for i, ch in enumerate(title_chars):
    add_textbox(slide, title_x, title_y + Inches(i * 0.52), Inches(0.55), Inches(0.55),
                ch, font_name=TITLE_FONT, font_size=44, color=HERO_RED, bold=True,
                alignment=PP_ALIGN.CENTER, line_spacing=1.0)

# 副标题 — 红绸上
add_textbox(slide, Inches(1.5), Inches(0.92), Inches(10), Inches(0.45),
            "2025-2026学年第二学期期中家长会",
            font_name=BODY_FONT, font_size=18, color=STAR_GOLD, bold=True,
            alignment=PP_ALIGN.CENTER, line_spacing=1.0)

# 底部信息
add_textbox(slide, Inches(8.0), Inches(6.7), Inches(4.5), Inches(0.4),
            "某某中学  |  七年级（3）班  |  2026.5.25",
            font_name=BODY_FONT, font_size=12, color=WHITE, bold=False,
            alignment=PP_ALIGN.RIGHT, line_spacing=1.0)


# ═══════════════════ 第2页：会议议程（行军路线） ═══════════════════
slide = prs.slides.add_slide(blank)
set_bg(slide, SMOKE_GRAY)
add_texture_tile(slide, manu_path)
# 叠加等高线纹理 — 简化
add_texture_tile(slide, rays_path)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.5),
            "行军路线", font_name=TITLE_FONT, font_size=28, color=HERO_RED, bold=True)

# 行军路线虚线
route_pts = [(Inches(1.5), Inches(1.2)), (Inches(3.5), Inches(2.0)),
             (Inches(5.5), Inches(3.0)), (Inches(8.0), Inches(4.5)), (Inches(10.5), Inches(5.8))]
for i in range(len(route_pts) - 1):
    add_line(slide, route_pts[i][0], route_pts[i][1], route_pts[i+1][0], route_pts[i+1][1],
             CAST_IRON, Pt(1.5))

# 四站
stations = [
    ("烽火台", "班情·烽火传讯", "约10分钟", Inches(1.5), Inches(1.2), False),
    ("号角", "学情·号角研判", "约15分钟", Inches(3.5), Inches(2.0), True),
    ("旗帜", "共育·旗帜所向", "约10分钟", Inches(5.5), Inches(3.0), False),
    ("红星", "对话·星光汇合", "约15分钟", Inches(8.0), Inches(4.5), False),
]

station_icons = {
    "烽火台": "🏰", "号角": "📯", "旗帜": "🚩", "红星": "⭐"
}

for name, label, duration, sx, sy, is_current in stations:
    # 图标圆形容器
    icon_bg = add_circle(slide, sx - Inches(0.15), sy - Inches(0.15), Inches(0.45),
                         fill=STAR_GOLD if is_current else SMOKE_GRAY,
                         line=STAR_GOLD if is_current else CAST_IRON, lw=Pt(2))
    if is_current:
        add_shadow(icon_bg, blur_pt=8, dist_pt=2, angle=135, opacity_pct=30)
        # 光芒射线
        glow_ring = add_circle(slide, sx - Inches(0.3), sy - Inches(0.3), Inches(0.75),
                               fill=STAR_GOLD)
        add_radial_glow(glow_ring, STAR_GOLD, center_alpha=30, edge_alpha=0)

    # 图标文字
    add_textbox(slide, sx - Inches(0.05), sy - Inches(0.05), Inches(0.25), Inches(0.25),
                station_icons.get(name, "●"), font_name=BODY_FONT, font_size=12,
                color=CAST_IRON, bold=False, alignment=PP_ALIGN.CENTER)

    # 站名
    add_textbox(slide, sx + Inches(0.4), sy - Inches(0.1), Inches(3.0), Inches(0.3),
                label, font_name=TITLE_FONT, font_size=18, color=HERO_RED, bold=True)
    # 时长
    add_textbox(slide, sx + Inches(0.4), sy + Inches(0.2), Inches(2.0), Inches(0.2),
                duration, font_name=BODY_FONT, font_size=11, color=GRAY_TEXT, bold=False)


# ═══════════════════ 第3页：学情透视（战旗高地） ═══════════════════
slide = prs.slides.add_slide(blank)
set_bg(slide, SMOKE_GRAY)
add_texture_tile(slide, manu_path)

# 左侧红旗背景
red_flag = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1.2), SLIDE_H,
                     fill=HERO_RED)
add_linear_gradient(red_flag, HERO_RED, CAST_IRON, angle=0)
# 旗面暗纹
add_textbox(slide, Inches(0.1), Inches(0.5), Inches(1.0), Inches(6.5),
            "中\n期\n战\n役\n数\n据", font_name=TITLE_FONT, font_size=16,
            color=STAR_GOLD, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=2.5)

# 标题
add_textbox(slide, Inches(1.8), Inches(0.4), Inches(8), Inches(0.5),
            "学情透视：期中战役战报", font_name=TITLE_FONT, font_size=26,
            color=HERO_RED, bold=True)

# 高地剪影轮廓
highland = add_shape(slide, MSO_SHAPE.WAVE, Inches(2.0), Inches(4.0), Inches(10.0), Inches(3.5),
                     fill=CAST_IRON)

# 战旗数据 — 用旗形表示
subjects = [
    ("语文", 82, 80, True),   # name, class_avg, grade_avg, advantage
    ("数学", 78, 83, False),
    ("英语", 71, 76, False),
    ("物理", 85, 80, True),
    ("历史", 88, 86, True),
]

flag_base_y = Inches(3.5)
for i, (subj, ca, ga, adv) in enumerate(subjects):
    fx = Inches(2.8) + Inches(1.9) * i
    # 旗杆高度映射分数
    class_h = Inches((ca - 50) / 50 * 2.2)
    grade_h = Inches((ga - 50) / 50 * 2.2)

    if adv:
        # 英雄红战旗 — 向上
        flag = add_shape(slide, MSO_SHAPE.RECTANGLE, fx, flag_base_y - class_h, Inches(0.6), class_h,
                         fill=HERO_RED)
        add_linear_gradient(flag, HERO_RED, STAR_GOLD, angle=90)
        # 金色飘带标记
        add_circle(slide, fx + Inches(0.15), flag_base_y - class_h - Inches(0.1), Inches(0.3),
                   fill=STAR_GOLD)
    else:
        # 铸铁色旗杆+下斜旗面
        flag_pole = add_shape(slide, MSO_SHAPE.RECTANGLE, fx + Inches(0.2), flag_base_y - class_h,
                              Inches(0.08), class_h, fill=CAST_IRON)
        flag = add_shape(slide, MSO_SHAPE.RECTANGLE, fx, flag_base_y - Inches(0.3),
                         Inches(0.5), Inches(0.25), fill=HERO_RED)
        # 警示标记
        warn = add_circle(slide, fx + Inches(0.1), flag_base_y - class_h - Inches(0.2),
                          Inches(0.25), fill=HERO_RED)
        add_textbox(slide, fx + Inches(0.12), flag_base_y - class_h - Inches(0.18),
                    Inches(0.2), Inches(0.2), "!", font_name=BODY_FONT, font_size=10,
                    color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 年级参考灰旗
    ref_flag = add_shape(slide, MSO_SHAPE.RECTANGLE, fx + Inches(0.9), flag_base_y - grade_h,
                         Inches(0.4), grade_h, fill=VICTORY_BLUE)

    # 科目名
    add_textbox(slide, fx, flag_base_y + Inches(0.15), Inches(1.2), Inches(0.25),
                subj, font_name=BODY_FONT, font_size=12, color=CAST_IRON, bold=True,
                alignment=PP_ALIGN.CENTER)

# 战地手记
add_textbox(slide, Inches(1.8), Inches(5.8), Inches(10.0), Inches(0.4),
            "▇ 英语防线巩固有效（+2分）      ▇ 数学阵地需重点支援（-5分）      ▇ 物理、历史优势保持",
            font_name=BODY_FONT, font_size=14, color=HERO_RED, bold=True)
add_textbox(slide, Inches(1.8), Inches(6.3), Inches(10.0), Inches(0.3),
            "战地手记：薄弱科目旗杆为铸铁色带警示标记，优势科目扬起金色飘带。",
            font_name=BODY_FONT, font_size=11, color=GRAY_TEXT, bold=False)


# ═══════════════════ 第4页：共育阵线（作战地图+弹药箱） ═══════════════════
slide = prs.slides.add_slide(blank)
set_bg(slide, SMOKE_GRAY)
add_texture_tile(slide, manu_path)
add_texture_tile(slide, rays_path)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
            "共育阵线：并肩作战", font_name=TITLE_FONT, font_size=26, color=HERO_RED, bold=True)

# 作战地图底
map_bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(7.5), Inches(5.5),
                   fill=SMOKE_GRAY, line=CAST_IRON, lw=Pt(1.5))
add_shadow(map_bg, blur_pt=6, dist_pt=2, angle=135, opacity_pct=10)

battles = [
    ("战场一：手机依赖", "⚡", "每日契约执行", Inches(1.2), Inches(1.5)),
    ("战场二：作业应付", "🔗", "检查思考痕迹", Inches(1.2), Inches(2.8)),
    ("战场三：情绪管理", "🛡", "每周15分钟倾听", Inches(1.2), Inches(4.1)),
]

for (title, icon, action, bx, by) in battles:
    # 红色箭头圈定
    arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, bx, by, Inches(0.4), Inches(0.3),
                      fill=HERO_RED)
    add_textbox(slide, bx + Inches(0.5), by - Inches(0.05), Inches(5.0), Inches(0.35),
                title, font_name=TITLE_FONT, font_size=16, color=HERO_RED, bold=True)
    # 图标
    add_textbox(slide, bx + Inches(5.5), by - Inches(0.05), Inches(0.4), Inches(0.35),
                icon, font_name=BODY_FONT, font_size=18, color=CAST_IRON, bold=False,
                alignment=PP_ALIGN.CENTER)
    # 行动文字
    add_textbox(slide, bx + Inches(0.5), by + Inches(0.3), Inches(5.5), Inches(0.25),
                action, font_name=BODY_FONT, font_size=13, color=STAR_GOLD, bold=True)

# 弹药箱
ammo_x = Inches(9.0); ammo_y = Inches(2.0)
ammo_box = add_shape(slide, MSO_SHAPE.RECTANGLE, ammo_x, ammo_y, Inches(3.5), Inches(4.0),
                     fill=CAST_IRON)
add_linear_gradient(ammo_box, CAST_IRON, RGBColor(0x5C, 0x3A, 0x1E), angle=0)
add_shadow(ammo_box, blur_pt=10, dist_pt=4, angle=135, opacity_pct=20)

# 弹药箱盖子
add_textbox(slide, ammo_x + Inches(0.3), ammo_y + Inches(0.2), Inches(3.0), Inches(0.6),
            "后方稳固，前线必胜", font_name=TITLE_FONT, font_size=18, color=HERO_RED,
            bold=True, alignment=PP_ALIGN.CENTER)

# 三个信封
envelopes = ["家庭行动\n锦囊①", "家庭行动\n锦囊②", "家庭行动\n锦囊③"]
for i, env_text in enumerate(envelopes):
    ey = ammo_y + Inches(1.1) + Inches(0.9) * i
    env = add_shape(slide, MSO_SHAPE.RECTANGLE, ammo_x + Inches(0.5), ey,
                    Inches(2.5), Inches(0.7), fill=SMOKE_GRAY, line=STAR_GOLD, lw=Pt(1))
    add_textbox(slide, ammo_x + Inches(0.6), ey + Inches(0.05), Inches(2.3), Inches(0.6),
                env_text, font_name=BODY_FONT, font_size=11, color=CAST_IRON, bold=True,
                alignment=PP_ALIGN.CENTER, line_spacing=1.4)


# ═══════════════════ 第5页：山河回响（和平与传承） ═══════════════════
slide = prs.slides.add_slide(blank)
set_bg(slide, SMOKE_GRAY)

# 天空渐变背景（硝烟暖灰→浅蓝）
sky = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(5.5),
                fill=SMOKE_GRAY)
add_linear_gradient(sky, SMOKE_GRAY, VICTORY_BLUE, angle=90)

# 白鸽群 — 放置鸽子剪影
dove_positions = [
    (Inches(10.0), Inches(3.5), 0.6, 20),
    (Inches(8.5), Inches(2.8), 0.8, 15),
    (Inches(7.0), Inches(2.0), 1.0, 10),
    (Inches(5.5), Inches(1.2), 0.7, 18),
    (Inches(4.0), Inches(0.6), 0.5, 25),
    (Inches(6.5), Inches(3.2), 0.5, 22),
    (Inches(9.5), Inches(1.5), 0.6, 12),
]
for dx, dy, scale, rot in dove_positions:
    dove = slide.shapes.add_picture(dove_path, dx, dy, Inches(scale), Inches(scale * 0.7))
    dove.rotation = rot

# 地平线
add_line(slide, Inches(0), Inches(4.8), SLIDE_W, Inches(4.8), STAR_GOLD, Pt(1.5))

# 少年剪影 — 简单矩形表示
for i in range(5):
    sx = Inches(10.0) + Inches(i * 0.4)
    # 身体
    add_shape(slide, MSO_SHAPE.RECTANGLE, sx, Inches(4.5), Inches(0.1), Inches(0.3),
              fill=CAST_IRON)
    # 头
    add_circle(slide, sx - Inches(0.02), Inches(4.35), Inches(0.12), fill=CAST_IRON)

# 收束语
add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10.0), Inches(0.8),
            "每一代人都有属于自己的长征，\n我们与孩子互为战友。",
            font_name=TITLE_FONT, font_size=26, color=HERO_RED, bold=True,
            alignment=PP_ALIGN.CENTER, line_spacing=1.8)

# 磨砂玻璃底条
frosted = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0),
                    fill=WHITE)
spPr_f = get_spPr(frosted)
solid_f = spPr_f.find(qn("a:solidFill"))
if solid_f is not None:
    srgb_f = solid_f.find(qn("a:srgbClr"))
    if srgb_f is not None:
        alpha_f = etree.SubElement(srgb_f, qn("a:alpha"))
        alpha_f.set("val", "70000")  # 70% opacity → 30% transparent

add_textbox(slide, Inches(1.0), Inches(6.6), Inches(5.0), Inches(0.3),
            "班主任电话：138-xxxx-xxxx  |  办公时间：周一至周五 8:00-17:00",
            font_name=BODY_FONT, font_size=12, color=CAST_IRON, bold=False,
            alignment=PP_ALIGN.LEFT, line_spacing=1.0)


# ═══════════════════ 保存 ═══════════════════
output = "E:/git/test_app/英雄主题家长会.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
